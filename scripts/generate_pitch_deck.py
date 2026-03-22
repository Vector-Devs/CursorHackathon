#!/usr/bin/env python3
"""
Generate a pitch deck (.pptx) styled like frontend/src/index.css (Riscon dashboard).
Run: ./scripts/build_pitch_deck.sh   (installs python-pptx if needed)
  or: pip install python-pptx && python3 scripts/generate_pitch_deck.py
Output: pitch-deck/CursorHackathon-Pitch.pptx

Slide flow: cover → problem → Riscon → revenue/GTM → team (why we) → why now → thank you (7 slides total).

Designed for ~4 minutes: cover (~25s) + 5 content slides (~35–40s each) + closing (~20s) — tighten talk track.
Each slide is one idea per bullet — expand verbally, not on the slide.

Detailed SAM/TAM sourcing lived in earlier deck iterations; keep footnotes in appendix if investors ask.

Market sizing notes (web research — verify definitions before investor meetings; syndicated reports mix software+services):
- SAM — EU: Statista statistic 1253133: 59,961 medium-sized (50–249 employees) manufacturing enterprises in the EU, 2024
  (largest sector share for that size band). Context: Eurostat Statistics Explained “Businesses in the manufacturing sector”:
  ~2.2 million manufacturing enterprises in the EU, 2023 (most are micro/small; SBS size-class tables: sbs_sc_ovw / sbs_ovw_act).
- SAM — US: U.S. Census Bureau County Business Patterns (CBP) or Economic Census NAICS 31–33, establishments by employment size
  (e.g. ECNLOCMFG); no single number substituted here—pull current national totals from data.census.gov.
- TAM — control tower (global, often software+services): WiseGuyReports summary (republished as Industry Today IT news): ~US$6.79B (2024),
  ~US$7.19B (2025) → ~US$12.7B (2035), ~5.9% CAGR. Narrower software-only forecasts: e.g. 360i Research / Research and Markets
  summaries cite ~US$3.94B (2025) → ~US$7.84B (2032) for control-tower software.
- TAM — supply chain risk management (SCRM) software: vendor reports vary widely (e.g. ~US$2.1–6.7B base year to ~US$26–46B
  2032–34 in trade-press summaries of SkyQuest, Verified Market Research, etc.); scope includes supplier risk, procurement, services.
- LLM unit cost: public API pricing pages; actual spend depends on article volume.

Slide 4 “Why now” (geopolitics): Middle East / Iran tensions and Gulf maritime risk are widely covered as affecting
shipping, insurance, and energy/trade flows (e.g. Strait of Hormuz, war-risk premiums). Cite major outlets (Reuters,
CNBC, etc.) for investor diligence — situation evolves; verify latest facts before presenting.
"""

import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.util import Inches, Pt

# --- Hackathon pop theme (high contrast; readable on projectors) ---
BG_PAGE = RGBColor(0x06, 0x08, 0x12)
BG_PAGE_ALT = RGBColor(0x0D, 0x11, 0x22)
BG_CARD = RGBColor(0x12, 0x18, 0x2A)
ACCENT = RGBColor(0x22, 0xE3, 0xFF)  # electric cyan
ACCENT_HOT = RGBColor(0xFF, 0x2D, 0x95)  # magenta punch
POP_LIME = RGBColor(0xCC, 0xFF, 0x00)  # acid lime
POP_VIOLET = RGBColor(0xA8, 0x5C, 0xFF)
TEXT_PRIMARY = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_SECONDARY = RGBColor(0xC8, 0xD4, 0xE8)
TEXT_TERTIARY = RGBColor(0x7A, 0x88, 0xA3)
SUCCESS = RGBColor(0x39, 0xFF, 0x9C)
WARNING = RGBColor(0xFF, 0xD0, 0x23)
BORDER = RGBColor(0x2A, 0x3A, 0x58)
BORDER_GLOW = RGBColor(0x22, 0xE3, 0xFF)

W = Inches(13.333)  # 16:9
H = Inches(7.5)
M = Inches(0.55)
TOP_BAR = Inches(0.18)


def add_full_bleed_bg(slide, color: RGBColor) -> None:
    """First shape added sits behind later shapes."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_hackathon_bg_extras(slide, *, cover: bool = False) -> None:
    """Geometric accents (after base bg). cover=True adds bigger drama for title slide."""
    # Skewed panel (depth)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.0), Inches(-0.6), Inches(7.2), Inches(9.2))
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor(0x08, 0x10, 0x22)
    band.line.fill.background()
    band.rotation = 11

    # Vertical neon rails
    for x_off, col in [(Inches(8.35), ACCENT_HOT), (Inches(8.52), ACCENT)]:
        rail = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_off, Inches(0.25), Inches(0.06), H - Inches(0.5))
        rail.fill.solid()
        rail.fill.fore_color.rgb = col
        rail.line.fill.background()

    if cover:
        # Big rings (fill matches bg so only stroke shows)
        o = slide.shapes.add_shape(MSO_SHAPE.OVAL, W - Inches(3.0), -Inches(0.5), Inches(4.0), Inches(4.0))
        o.fill.solid()
        o.fill.fore_color.rgb = BG_PAGE
        o.line.color.rgb = POP_LIME
        o.line.width = Pt(5)

        o2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, -Inches(1.5), H - Inches(2.8), Inches(4.2), Inches(3.2))
        o2.fill.solid()
        o2.fill.fore_color.rgb = BG_PAGE
        o2.line.color.rgb = ACCENT
        o2.line.width = Pt(4)


def add_accent_bar(slide, *, hot_tip: bool = False) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, TOP_BAR)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_HOT if hot_tip else ACCENT
    bar.line.fill.background()
    tip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, W - Inches(2.8), 0, Inches(2.8), TOP_BAR)
    tip.fill.solid()
    tip.fill.fore_color.rgb = POP_LIME
    tip.line.fill.background()


def add_card(slide, left, top, width, height, *, glow: bool = True) -> None:
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = BG_CARD
    if glow:
        card.line.color.rgb = BORDER_GLOW
        card.line.width = Pt(2.5)
    else:
        card.line.color.rgb = BORDER
        card.line.width = Pt(1)


def add_left_slide_accent(slide, top_inches: float, height_inches: float) -> None:
    stripe = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, M, Inches(top_inches), Inches(0.12), Inches(height_inches)
    )
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = ACCENT_HOT
    stripe.line.fill.background()


def apply_textframe_bullets(tf) -> None:
    """Turn each paragraph into a real PPTX bullet (•)."""
    for p in tf.paragraphs:
        pPr = p._p.get_or_add_pPr()
        bu = parse_xml(
            '<a:buChar xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" char="•"/>'
        )
        pPr.append(bu)


def _add_md_runs_to_paragraph(p, line: str, *, font_pt: int) -> None:
    if len(p.runs) == 0:
        p.add_run()
    segments = re.split(r"(\*\*[^*]+\*\*)", line)
    first_seg = True
    for seg in segments:
        if not seg:
            continue
        if first_seg:
            run = p.runs[0]
            first_seg = False
        else:
            run = p.add_run()
        if seg.startswith("**") and seg.endswith("**"):
            run.text = seg[2:-2]
            run.font.bold = True
            run.font.size = Pt(font_pt)
            run.font.color.rgb = ACCENT
            run.font.name = "Calibri"
        else:
            run.text = seg
            run.font.bold = False
            run.font.size = Pt(font_pt)
            run.font.color.rgb = TEXT_SECONDARY
            run.font.name = "Calibri"


def populate_tf_markdown(tf, lines: list[str], body_pt: int = 15) -> None:
    """Fill text frame; **bold** becomes accent-colored bold (markdown-style)."""
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        p.line_spacing = 1.12
        _add_md_runs_to_paragraph(p, line, font_pt=body_pt)


def populate_tf_markdown_levels(tf, items: list[tuple[str, int]], body_pt: int = 15) -> None:
    """Like populate_tf_markdown, with outline levels 0 (main), 1 (sub), 2 (sub-sub), …"""
    for i, (line, level) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = min(level, 8)
        sub = level > 0
        p.space_after = Pt(4 if level > 1 else (6 if sub else 12))
        p.line_spacing = 1.1
        font_pt = max(10, body_pt - 2 * min(level, 3))
        if len(p.runs) == 0:
            p.add_run()
        _add_md_runs_to_paragraph(p, line, font_pt=font_pt)


def textbox(slide, left, top, width, height, text: str, *, size=14, bold=False, color=TEXT_PRIMARY, font="Calibri", align=PP_ALIGN.LEFT, small_caps=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    r = p.runs[0]
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    if small_caps:
        try:
            r.font.small_caps = True
        except AttributeError:
            pass
    return box


def bullet_slide(
    prs,
    title: str,
    lines: list[str],
    kicker: str | None = None,
    *,
    body_pt: int = 14,
    ppt_bullets: bool = False,
    lines_with_levels: list[tuple[str, int]] | None = None,
):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bleed_bg(slide, BG_PAGE)
    add_hackathon_bg_extras(slide, cover=False)
    add_accent_bar(slide)

    if kicker:
        textbox(
            slide,
            M,
            Inches(0.38),
            W - 2 * M,
            Inches(0.4),
            kicker.upper(),
            size=11,
            color=POP_LIME,
            bold=True,
            font="Consolas",
            small_caps=True,
        )

    textbox(
        slide,
        M,
        Inches(0.78),
        W - 2 * M,
        Inches(1.0),
        title.upper(),
        size=34,
        bold=True,
        color=TEXT_PRIMARY,
        font="Calibri Light",
    )

    body_top = Inches(1.88)
    body_h = H - body_top - M
    _top_in = 1.88
    _h_in = 7.5 - 0.55 - _top_in
    add_left_slide_accent(slide, _top_in, _h_in)
    add_card(slide, M, body_top, W - 2 * M, body_h, glow=True)

    tb = slide.shapes.add_textbox(M + Inches(0.42), body_top + Inches(0.38), W - 2 * M - Inches(0.78), body_h - Inches(0.6))
    tf = tb.text_frame
    tf.word_wrap = True
    if lines_with_levels is not None:
        populate_tf_markdown_levels(tf, lines_with_levels, body_pt=body_pt)
    else:
        populate_tf_markdown(tf, lines, body_pt=body_pt)
    if ppt_bullets:
        apply_textframe_bullets(tf)
    return slide


def title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bleed_bg(slide, BG_PAGE)
    add_hackathon_bg_extras(slide, cover=True)
    add_accent_bar(slide)

    # HACKATHON badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, M, Inches(0.55), Inches(2.55), Inches(0.52))
    badge.fill.solid()
    badge.fill.fore_color.rgb = ACCENT_HOT
    badge.line.fill.background()
    bt = slide.shapes.add_textbox(M + Inches(0.2), Inches(0.58), Inches(2.2), Inches(0.45))
    btf = bt.text_frame
    btf.paragraphs[0].text = "⚡ HACKATHON DEMO"
    btf.paragraphs[0].font.size = Pt(11)
    btf.paragraphs[0].font.bold = True
    btf.paragraphs[0].font.name = "Consolas"
    btf.paragraphs[0].font.color.rgb = TEXT_PRIMARY

    # Split headline color for punch (RIS + CON)
    textbox(
        slide,
        M,
        Inches(1.35),
        Inches(1.45),
        Inches(1.2),
        "RIS",
        size=72,
        bold=True,
        color=TEXT_PRIMARY,
        font="Calibri Light",
        align=PP_ALIGN.LEFT,
    )
    textbox(
        slide,
        M + Inches(1.38),
        Inches(1.35),
        Inches(3.2),
        Inches(1.2),
        "CON",
        size=72,
        bold=True,
        color=ACCENT,
        font="Calibri Light",
        align=PP_ALIGN.LEFT,
    )

    textbox(
        slide,
        M,
        Inches(2.58),
        W - 2 * M,
        Inches(0.55),
        "SUPPLY CHAIN RISK COMMAND CENTER",
        size=15,
        bold=True,
        color=POP_LIME,
        font="Consolas",
        align=PP_ALIGN.LEFT,
        small_caps=True,
    )
    textbox(
        slide,
        M,
        Inches(3.25),
        W - 2 * M,
        Inches(1.15),
        "Maritime & logistics signals → live probability → your footprint on a map.",
        size=20,
        color=TEXT_SECONDARY,
        font="Calibri",
        align=PP_ALIGN.LEFT,
    )

    chip_y = Inches(5.05)
    chip_w = Inches(2.92)
    chip_h = Inches(0.72)
    gap = Inches(0.28)
    labels = [
        ("LIVE SIGNALS", SUCCESS),
        ("PROBABILITY", ACCENT),
        ("ENTERPRISE DATA", ACCENT_HOT),
    ]
    x = M
    for label, c in labels:
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, chip_y, chip_w, chip_h)
        sh.fill.solid()
        sh.fill.fore_color.rgb = RGBColor(0x0C, 0x12, 0x22)
        sh.line.color.rgb = c
        sh.line.width = Pt(2)
        t = slide.shapes.add_textbox(x + Inches(0.22), chip_y + Inches(0.16), chip_w - Inches(0.4), chip_h)
        tf = t.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.name = "Consolas"
        p.font.color.rgb = c
        x += chip_w + gap

    textbox(
        slide,
        M,
        H - Inches(0.65),
        W - 2 * M,
        Inches(0.45),
        "CursorHackathon · Vector-Devs",
        size=12,
        color=TEXT_TERTIARY,
        font="Calibri",
    )
    return slide


def main():
    out_dir = Path(__file__).resolve().parent.parent / "pitch-deck"
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / "CursorHackathon-Pitch.pptx"

    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # Cover
    title_slide(prs)

    # Slide 1 — Problem (~40s): concise bullets + real PPTX bullets; facts under point 1 (researched)
    bullet_slide(
        prs,
        "The problem",
        [],
        kicker="1 / 5",
        body_pt=15,
        ppt_bullets=True,
        lines_with_levels=[
            (
                "**No live link** — alerts & spreadsheets aren’t tied to plants, suppliers, or in‑transit goods.",
                0,
            ),
            (
                "EU: **~59,961** medium-sized manufacturers (50–249 employees), **2024** — Statista (stat. 1253133).",
                1,
            ),
            (
                "EU manufacturing: **~2.2 million** enterprises total (**2023**) — **Eurostat** (most are micro/small; mid-market = ICP).",
                1,
            ),
            (
                "Global **control-tower** market often quoted **~US$6.8–7.2B (2024–25)** → **~US$12.7B by 2035** — analyst summaries (e.g. WiseGuyReports / Industry Today).",
                1,
            ),
            (
                "**How it’s tackled today** — email chains, shared **Excel** risk registers, and siloed spreadsheets; **no** single live view across plants & suppliers.",
                0,
            ),
            (
                "**~24–48h** typical lag to align on exposure & action — manual handoffs, not live signals. **Kinaxis** research (**2024**): **83%** of supply chains **cannot** respond to disruptions within **24 hours**.",
                1,
            ),
            (
                "**Large EU+US pool** — tens of thousands of addressable mid-market firms; multi‑$B risk / visibility software TAMs.",
                0,
            ),
            (
                "**BMW** — **Connected Supply Chain Cockpit** for 40+ plants & real-time material visibility (public implementation case studies).",
                1,
            ),
            (
                "**Siemens** — **supply-chain control-tower** / end-to-end visibility software & white papers (industrial / CP&R).",
                1,
            ),
            (
                "**P&G**, **Merck**, **Castrol (BP)** — named **Kinaxis** customers for large-scale supply-chain planning & transparency (vendor case studies / press).",
                1,
            ),
        ],
    )

    # Slide 2 — The solution — Riscon: wedge + definition + grouped customer gains (nested bullets)
    bullet_slide(
        prs,
        "The solution — Riscon",
        [],
        kicker="2 / 4",
        body_pt=14,
        ppt_bullets=True,
        lines_with_levels=[
            ("**Our wedge** — **headline noise → portfolio-level risk** on one live board...", 0),
            (
                "**What Riscon is** — **Live supply-chain risk cockpit:** AI **classifies maritime & logistics news in real time**, **fuses** it with **your** plants, suppliers, and **fleet context**, and **streams probability** to the control-tower UI (**WebSockets** — not weekly Excel packs).",
                0,
            ),
            ("**Customer gains** — why buyers care:", 0),
            (
                "**Speed** — **same-day**, **site-aware** signals vs. **24–48h** email / spreadsheet coordination.",
                1,
            ),
            (
                "**Clarity** — **which plants & suppliers** are in the risk zone, not a **generic** news feed.",
                1,
            ),
            (
                "**Alignment** — **one map + dashboard** for **ops, procurement, leadership** — fewer crisis calls, faster decisions.",
                1,
            ),
            (
                "**Proof in the demo** — **news → probability → map** in one screen · **Docker** microservices (Spring, React, agents) for **fast pilots**.",
                0,
            ),
        ],
    )

    # Slide 3 — Business + GTM: grouped bullets
    bullet_slide(
        prs,
        "Revenue · go-to-market",
        [],
        kicker="3 / 5",
        body_pt=14,
        ppt_bullets=True,
        lines_with_levels=[
            ("**Who buys**", 0),
            (
                "VP Supply Chain / CPO / Head of Logistics; **3PLs / 4PLs** later for **white-label** risk views to shippers.",
                1,
            ),
            ("**Pricing** — anchors (tune to ICP)", 0),
            (
                "**~$25–75/user/mo** team tier **or** **~$18k–60k/yr** platform for mid-market (10–50 sites) + minimums.",
                1,
            ),
            ("**Economics** (illustrative)", 0),
            (
                "**~$40k ACV**, **~75%** gross margin after cloud + LLM; **cap** AI usage; **<12 mo** CAC payback via inside sales + pilots.",
                1,
            ),
            ("**First wins** — GTM", 0),
            (
                "**Manufacturers** — intros via **VDMA / IHK / trade shows**; **30-day pilots**.",
                1,
            ),
            ("**3PL** partners — co-branded **shipper risk** widget.", 1),
            ("**One public case** — “alert → action in minutes.” [EDIT: your distribution].", 1),
        ],
    )

    # Slide 4 — Why we (team only)
    bullet_slide(
        prs,
        "Team — Vector Devs",
        [],
        kicker="4 / 5",
        body_pt=14,
        ppt_bullets=True,
        lines_with_levels=[
            ("**Why we** — execution & team", 0),
            (
                "**Vector Devs** — **3 developers**, **dream team**: **deep experience**, **rapid prototyping**, end-to-end ownership.",
                1,
            ),
            (
                "**Hackathon build:** full stack — agents, probability service, React dashboard, Docker — **delivered in 27.5 hours** (not slides only — **shipped**).",
                1,
            ),
            ("**Meet the team** (LinkedIn)", 1),
            (
                "**Balaji Vengatesh M** — https://www.linkedin.com/in/balaji-vengatesh-m-3b871151/",
                2,
            ),
            (
                "**Sven Schmeckenbecher** — https://www.linkedin.com/in/sven-schmeckenbecher",
                2,
            ),
            (
                "**Sreekanth Kata** — https://www.linkedin.com/in/sreekanth-kata-3096471bb/",
                2,
            ),
        ],
    )

    # Slide 5 — Why now (market + geopolitical urgency)
    bullet_slide(
        prs,
        "Why now",
        [],
        kicker="5 / 5",
        body_pt=14,
        ppt_bullets=True,
        lines_with_levels=[
            ("**Market + geopolitical urgency**", 0),
            (
                "**Middle East / Iran tensions** — escalations have **raised maritime supply-chain risk** (e.g. **Strait of Hormuz** chokepoint, **war-risk insurance** & **freight** repricing, **carrier rerouting**). **Global trade & energy** flows get **scrutiny**; coverage in **Reuters**, **CNBC**, **Foreign Policy**, **Janes** (2025–26) shows **shipping** and **boards** treating disruption as **immediate**, not theoretical.",
                1,
            ),
            (
                "**Customer sensitivity** — teams need **live** visibility into **suppliers, lanes, and sites** when **geopolitical shocks** hit; **Excel + email** cannot keep pace. **Riscon** maps news → **probability** → **your** footprint.",
                1,
            ),
            (
                "**Tech window** — **LLM APIs** make real-time classification affordable; **resilience** budgets and **board** pressure are up; incumbents still aren’t **map-first** — room for a **focused** control tower.",
                1,
            ),
        ],
    )

    # Thank you
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bleed_bg(slide, BG_PAGE)
    add_hackathon_bg_extras(slide, cover=False)
    add_accent_bar(slide, hot_tip=True)
    textbox(
        slide,
        M,
        Inches(2.35),
        W - 2 * M,
        Inches(1.35),
        "THANK YOU",
        size=52,
        bold=True,
        color=TEXT_PRIMARY,
        font="Calibri Light",
        align=PP_ALIGN.CENTER,
    )
    textbox(
        slide,
        M,
        Inches(3.85),
        W - 2 * M,
        Inches(0.65),
        "QUESTIONS?",
        size=26,
        color=POP_LIME,
        font="Consolas",
        align=PP_ALIGN.CENTER,
        bold=True,
    )
    textbox(
        slide,
        M,
        Inches(4.75),
        W - 2 * M,
        Inches(0.85),
        "github.com/Vector-Devs/CursorHackathon",
        size=15,
        color=ACCENT,
        font="Consolas",
        align=PP_ALIGN.CENTER,
        bold=True,
    )

    prs.save(out_path)
    print(f"Wrote {out_path}")


if __name__ == "__main__":
    main()
